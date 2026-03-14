import os
import io
import requests
from bs4 import BeautifulSoup
from groq import Groq
from docx import Document
from urllib.parse import urlparse
from pptx import Presentation
import pypdf

# Groq model to use - llama-3.3-70b-versatile is free and high quality
GROQ_MODEL = 'llama-3.3-70b-versatile'

def extract_text_from_pptx(file_bytes):
    """Extracts all text from an uploaded PPTX file (as bytes)."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
        return "\n".join(text_runs)
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ""

def extract_text_from_pdf(file_bytes):
    """Extracts all text from an uploaded PDF file (as bytes)."""
    try:
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""


PROMPTS = {
    'profile': """You are a professional business writer. Based on the following information extracted from a company's website, write a professional company profile.
The profile should be around 2-3 paragraphs. It should explain who the company is, what they do, their key objectives, and their macroeconomic/market impact (if apparent). Include information on their international trade or expansion efforts if available. Do not include introductory conversational text like 'Here is the profile...'.

Company Website Content:
{context}
""",
    'faq': """You are a professional business writer. Based on the following information extracted from a company's website, generate a "Frequently Asked Questions" (FAQ) document.
Create 4-5 relevant questions and answer them clearly and concisely based on the text. Examples of good questions are about their history, size of companies they work with, governments they work with, and regions they operate in. Do not use conversational filler.

Company Website Content:
{context}
""",
    'services': """You are a professional business writer. Based on the following information extracted from a company's website, generate a "Services Overview" document.
Group their services into 2-3 logical categories (e.g., 'Business Consulting', 'Global Expansion', etc.). For each category, provide a short paragraph or bullet points explaining what they offer. Do not include conversational filler in the output.

Company Website Content:
{context}
""",
    'success_stories': """You are a professional business writer. Based on the following information extracted from a company's website, generate a "Success Stories and Proof Points" document.
Extract any metrics, track records, case studies, or proof points from the text and present them as a clean list of achievements. If none are explicitly stated, formulate general potential value propositions they offer based on their services. Do not include conversational filler.

Company Website Content:
{context}
""",
    'objection_responses': """You are a professional sales strategist. Based on the following information extracted from a company's website, write "Objection Responses" for their sales team.
Create responses for the following 4 common objections, tailored to the company's specific services and value proposition:
1. "I’m not interested."
2. "Send me an email instead."
3. "Who are you and why are you calling?"
4. "We already have a consulting partner."

Provide a short, professional 2-3 sentence response for each objection. Do not include conversational filler in the output.

Company Website Content:
{context}
"""
}

def extract_text_from_url(url):
    """Scrapes all text from a given URL."""
    try:
        # Add a custom user agent to avoid basic blocks
        headers = {'User-Agent': 'Mozilla/5.0'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.extract()
            
        text = soup.get_text(separator=' ')
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text
    except Exception as e:
        print(f"Error scraping {url}: {e}")
        return ""

def generate_content(prompt_key, context, api_key=None):
    """Uses Groq cloud API to generate text based on the prompt template."""
    if not context:
        return "Not enough information to generate this document."

    prompt = PROMPTS[prompt_key].format(context=context[:4000])  # limit context to avoid overflowing context window

    # Use provided key or fall back to environment variable
    key = api_key or os.environ.get("GROQ_API_KEY", "")
    if not key:
        return "Error: GROQ_API_KEY not set. Please add it to your Streamlit secrets or environment."

    print(f"Generating '{prompt_key}' using Groq model {GROQ_MODEL}...")
    try:
        client = Groq(api_key=key)
        completion = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1024,
        )
        return completion.choices[0].message.content.strip()
    except Exception as e:
        print(f"Error calling Groq: {e}")
        return f"Error generation failed: {e}"

def save_to_docx(content, folder, filename, title):
    """Saves generated markdown/text to a .docx file."""
    doc = Document()
    doc.add_heading(title, 0)
    
    # Simple parsing: double newlines usually mean new paragraph
    paragraphs = content.split('\n\n')
    for p in paragraphs:
        if p.strip():
            doc.add_paragraph(p.strip())
            
    filepath = os.path.join(folder, f"{filename}.docx")
    doc.save(filepath)
    print(f"Saved {filepath}")

def process_organization(url):
    """End-to-end process for a single organization URL."""
    # Create valid directory name from URL
    domain = urlparse(url).netloc.replace("www.", "")
    safe_domain_name = domain.split('.')[0] # e.g., 'nmpartnership'
    
    # Create the folder if it doesn't exist
    if not os.path.exists(safe_domain_name):
        os.makedirs(safe_domain_name)
    
    print(f"\n==========================================")
    print(f"Processing Organization: {url}")
    print(f"==========================================")
    
    print("1. Scraping website context...")
    context = extract_text_from_url(url)
    
    if not context:
        print(f"Failed to extract context for {url}. Skipping docs generation.")
        return
        
    titles = {
        'profile': "Company Profile",
        'faq': "Frequently Asked Questions",
        'services': "Services Overview",
        'success_stories': "Success Stories and Proof Points",
        'objection_responses': "Objection Responses"
    }

    print("2. Generating documents via AI...")
    for key, title in titles.items():
        generated_text = generate_content(key, context)
        save_to_docx(generated_text, safe_domain_name, key, title)
        
    print(f"Successfully processed {url} -> Documents saved in './{safe_domain_name}/'")

if __name__ == "__main__":
    # Example lists given by user
    urls = [
        "https://nmpartnership.com/",
        "https://www.htai.de/"
    ]
    
    print("Starting Automated Documentation Generator")
    for u in urls:
        process_organization(u)
    print("\nOverall processing complete!")
